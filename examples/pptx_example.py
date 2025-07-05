"""Example usage of the PptxParser.

This script generates a small sample PowerPoint file on the fly (to avoid
committing binary assets) and then parses it using the library. The resulting
Markdown is printed to stdout.
"""

from __future__ import annotations

import asyncio
from pathlib import Path

from pptx import Presentation  # type: ignore
from pptx.util import Inches, Pt  # type: ignore

from doc_parser.core import ParserConfig, ParserRegistry
from doc_parser.utils import save_markdown


def _build_sample_pptx(path: Path) -> None:
    """Create a minimal two-slide PPTX for demonstration purposes."""
    prs = Presentation()

    # Slide 1 – title & bullet list
    slide_layout = prs.slide_layouts[1]  # Title & content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Document Parser Demo"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Key Features"
    for line in [
        "Supports PDF, DOCX, XLSX, PPTX, HTML",
        "AI-powered extraction",
        "Smart caching",
    ]:
        p = body.add_paragraph()
        p.text = line
        p.level = 1

    # Slide 2 – table example
    slide_layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Simple Table"

    rows, cols = 3, 3
    left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(1.5)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table

    # Header row formatting
    headers = ["Col A", "Col B", "Col C"]
    for col, text in enumerate(headers):
        cell = tbl.cell(0, col)
        cell.text = text
        cell.text_frame.paragraphs[0].font.bold = True

    # Data rows
    for row in range(1, rows):
        for col in range(cols):
            tbl.cell(row, col).text = f"R{row}C{col}"

    prs.save(path)


async def _demo() -> None:
    
    sample_path = Path("outputs/sample_demo.pptx")
    cfg = ParserConfig(output_format="markdown")
    parser = ParserRegistry.get_parser(sample_path, cfg)

    result = await parser.parse_with_cache(sample_path)

    # Save results to a Markdown file
    output_md = Path("outputs") / f"{sample_path.stem}_parsed.md"
    save_markdown(result.content, output_md)
    print(f"Parsed markdown saved to: {output_md}")


if __name__ == "__main__":
    asyncio.run(_demo()) 