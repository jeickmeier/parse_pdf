"""Unit tests for PptxParser."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation  # type: ignore
from pptx.util import Inches

from doc_parser.core.settings import Settings
from doc_parser.core.registry import ParserRegistry

# mypy: ignore-errors

# Mark entire module as asyncio-compatible, allowing removal of individual decorators
pytestmark = pytest.mark.asyncio


def _create_temp_pptx(tmp_path: Path) -> Path:
    pptx_path = tmp_path / "temp.pptx"
    prs = Presentation()

    # Slide with title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Slide"

    # Simple table
    rows, cols = 2, 2
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(1), Inches(2), Inches(6), Inches(1)
    )
    tbl = table_shape.table
    tbl.cell(0, 0).text = "Header A"
    tbl.cell(0, 1).text = "Header B"
    tbl.cell(1, 0).text = "Val 1"
    tbl.cell(1, 1).text = "Val 2"

    prs.save(pptx_path)
    return pptx_path


@pytest.mark.asyncio
async def test_pptx_parser_markdown(tmp_path: Path) -> None:
    pptx_path = _create_temp_pptx(tmp_path)

    cfg = Settings(output_format="markdown", use_cache=False)
    parser = ParserRegistry.from_path(pptx_path, cfg)

    result = await parser.parse(pptx_path)

    assert not result.errors
    assert "# Test Slide" in result.content
    assert "| Header A" in result.content  # table header


@pytest.mark.asyncio
async def test_pptx_parser_json(tmp_path: Path) -> None:
    pptx_path = _create_temp_pptx(tmp_path)

    cfg = Settings(output_format="json", use_cache=False)
    parser = ParserRegistry.from_path(pptx_path, cfg)

    result = await parser.parse(pptx_path)

    assert not result.errors
    import json as _json

    data = _json.loads(result.content)
    assert isinstance(data, list)
    assert data[0]["title"] == "Test Slide"
    assert any("Header A" in cell for row in data[0]["tables"][0] for cell in row)
