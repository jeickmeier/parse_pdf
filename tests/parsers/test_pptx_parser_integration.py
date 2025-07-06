import asyncio
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from doc_parser.config import AppConfig
from doc_parser.parsers.pptx.parser import PptxParser


@pytest.fixture()
def make_sample_pptx(tmp_path):  # noqa: D401
    """Return path to a generated PPTX with simple content."""

    def _factory() -> Path:
        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = "Sample Deck"
        slide.placeholders[1].text = "Subtitle"

        # Content slide with bullet list
        bullet_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_layout)
        body = slide2.shapes.placeholders[1].text_frame
        body.text = "First"
        body.add_paragraph().text = "Second"

        # Save
        path = tmp_path / "deck.pptx"
        prs.save(path)
        return path

    return _factory


@pytest.mark.asyncio
async def test_pptx_parser_markdown(make_sample_pptx):
    pptx_path = make_sample_pptx()

    settings = AppConfig(
        output_format="markdown",
        parser_settings={"pptx": {"extract_images": False, "slide_delimiter": "---"}},
    )
    parser = PptxParser(settings)

    result = await parser.parse(pptx_path)
    # Should contain title as markdown heading
    assert result.content.startswith("# Sample Deck")
    # bullet list captured
    assert "Second" in result.content
    # Slide delimiter present
    assert "---" in result.content
    assert result.metadata["slides"] == 2


@pytest.mark.asyncio
async def test_pptx_parser_json(make_sample_pptx):
    pptx_path = make_sample_pptx()
    settings = AppConfig(output_format="json", use_cache=False, parser_settings={"pptx": {"extract_images": False}})
    parser = PptxParser(settings)
    result = await parser.parse(pptx_path)
    assert "\"title\"" in result.content
    assert "\"text\"" in result.content


def test_pptx_validate_input_neg(tmp_path):
    fake = tmp_path / "file.txt"
    fake.write_text("x")
    parser = PptxParser(AppConfig())
    assert asyncio.run(parser.validate_input(fake)) is False 